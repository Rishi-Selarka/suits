#!/usr/bin/env python3
"""
Suits AI — Hackathon Presentation Generator v2
Theme: Dark (#111114) + Orange (#DA6B2B) + Cream (#FAF9EA) — Claude-inspired
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──────────────────────────────────────────────────────────────
BASE = "/Users/rishiselarka/Documents/suits"
SS_DIR = os.path.join(BASE, "pptss", "converted")
LOGO = os.path.join(BASE, "suits/frontend/public/images/suits-logo.png")
OUTPUT = os.path.join(BASE, "Suits_AI_Presentation.pptx")

# ── Theme Colors ───────────────────────────────────────────────────────
BG_DARK = RGBColor(0x11, 0x11, 0x14)
BG_DARKER = RGBColor(0x09, 0x09, 0x0B)
ORANGE = RGBColor(0xDA, 0x6B, 0x2B)
ORANGE_LIGHT = RGBColor(0xE0, 0x77, 0x33)
ORANGE_SOFT = RGBColor(0xF5, 0xAD, 0x7A)
CREAM = RGBColor(0xFA, 0xF9, 0xEA)
CREAM_DIM = RGBColor(0xA1, 0xA1, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE_200 = RGBColor(0x18, 0x18, 0x1B)
SURFACE_300 = RGBColor(0x27, 0x27, 0x2A)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE_ACCENT = RGBColor(0x60, 0xA5, 0xFA)

# ── Presentation Setup (16:9) ──────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# Screenshot aspect ratio: 3024 x 1722 ≈ 1.756:1
SS_RATIO = 1722 / 3024  # height / width = 0.5694


# ── Helper Functions ───────────────────────────────────────────────────
def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             color=CREAM, bold=False, align=PP_ALIGN.LEFT, name="Arial"):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return box


def add_rich_text(slide, left, top, width, height, runs_list):
    """runs_list = list of (text, size, color, bold) tuples per paragraph.
       Each item is one paragraph. Use \\n inside text for manual breaks."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, sz, clr, bld) in enumerate(runs_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Arial"
        p.space_after = Pt(4)
    return box


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, text="", font_size=12,
             font_color=CREAM, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_thin_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE, width=2):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def add_line(slide, x1, y1, x2, y2, color=SURFACE_300, width=1):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def add_img(slide, path, left, top, width):
    """Add image with auto-calculated height from SS_RATIO."""
    if not os.path.exists(path):
        return None
    h = int(width * SS_RATIO)
    return slide.shapes.add_picture(path, left, top, width, h)


def add_circle(slide, left, top, size, fill_color, text="", fs=11,
               fc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    if text:
        p = s.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = fc
        p.font.bold = True
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
    return s


def section_header(slide, label, title, subtitle=None):
    """Standard section header: orange label + line + big title + optional subtitle."""
    add_text(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.35),
             label, size=12, color=ORANGE, bold=True)
    add_thin_bar(slide, Inches(0.8), Inches(0.92), Inches(2.8), Inches(0.04), ORANGE)
    add_text(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.7),
             title, size=28, color=CREAM, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.85), Inches(11), Inches(0.6),
                 subtitle, size=14, color=CREAM_DIM)


def slide_footer(slide, left_text="RNSIT Agentic AI Hackathon  |  PS-3",
                 right_text="Suits AI"):
    """Subtle bottom bar."""
    add_thin_bar(slide, Inches(0), Inches(7.1), W, Inches(0.4), SURFACE_200)
    add_text(slide, Inches(0.8), Inches(7.13), Inches(5), Inches(0.35),
             left_text, size=9, color=CREAM_DIM)
    add_text(slide, Inches(8), Inches(7.13), Inches(4.5), Inches(0.35),
             right_text, size=9, color=ORANGE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

# Top accent
add_thin_bar(s, Inches(0), Inches(0), W, Inches(0.05), ORANGE)

# Logo centered
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.92), Inches(1.2), Inches(1.5))

# Title
add_text(s, Inches(1.5), Inches(2.9), Inches(10.333), Inches(1.0),
         "SUITS AI", size=64, color=CREAM, bold=True, align=PP_ALIGN.CENTER)

# Tagline
add_text(s, Inches(2), Inches(4.0), Inches(9.333), Inches(0.6),
         "Your Legal Intelligence, Tailored to You",
         size=22, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

# Subtitle
add_text(s, Inches(2.5), Inches(4.7), Inches(8.333), Inches(0.5),
         "Multi-Agent AI System for Legal Document Analysis & Decision Support",
         size=14, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Divider line
add_line(s, Inches(5.5), Inches(5.5), Inches(7.833), Inches(5.5), SURFACE_300, 1)

# Bottom context
add_text(s, Inches(2), Inches(5.7), Inches(9.333), Inches(0.4),
         "RNSIT Agentic AI Hackathon  |  Problem Statement 3  |  AI Legal Document Action Agent",
         size=11, color=CREAM_DIM, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "THE PROBLEM",
               "Everyone signs legal documents.\nVery few truly understand them.")

problems = [
    ("1", "Complex Language Barrier",
     "Rental agreements, NDAs, and terms of service are filled\n"
     "with legal jargon — creating a dangerous gap between\n"
     "what people sign and what they actually understand.",
     ORANGE),
    ("2", "Hidden Risks & Obligations",
     "Critical penalties, liability clauses, and auto-renewal traps\n"
     "are buried in dense paragraphs. Easy to miss, costly to overlook.",
     RED),
    ("3", "No Accessible Guidance",
     "Most people can't afford a lawyer for every document.\n"
     "They need a smart assistant that makes legal text approachable.",
     AMBER),
    ("4", "Uninformed Decision-Making",
     "Without understanding what they're agreeing to, people sign\n"
     "away rights, accept unfair terms, and expose themselves to risk.",
     BLUE_ACCENT),
]

for i, (num, title, desc, accent) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(2.7) + Inches(row * 2.15)

    # Card
    add_rect(s, x, y, Inches(5.8), Inches(1.85), SURFACE_200, border_color=SURFACE_300)
    # Left accent stripe
    add_thin_bar(s, x, y, Inches(0.06), Inches(1.85), accent)
    # Number circle
    add_circle(s, x + Inches(0.3), y + Inches(0.25), Inches(0.4), accent, num, fs=15)
    # Title
    add_text(s, x + Inches(0.9), y + Inches(0.2), Inches(4.5), Inches(0.35),
             title, size=16, color=CREAM, bold=True)
    # Desc
    add_text(s, x + Inches(0.9), y + Inches(0.6), Inches(4.6), Inches(1.1),
             desc, size=11, color=CREAM_DIM)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "OUR SOLUTION",
               "Suits AI — A Multi-Agent Legal Intelligence Platform",
               "6 specialized AI agents work together to analyze, simplify, assess risk,\n"
               "benchmark, advise, and verify — so you can make informed decisions.")

features = [
    ("Plain Language\nSimplifier", "Breaks down complex legal\ntext into simple English\nanyone can understand.", ORANGE),
    ("Clause-Level\nRisk Scoring", "Scores every clause 1-10,\nflags dangerous patterns\nwith color-coded alerts.", RED),
    ("Contextual\nQ&A via RAG", "Ask questions about your\ndocument — get clear,\nsource-grounded answers.", BLUE_ACCENT),
    ("Industry\nBenchmark", "Compares your clauses\nagainst fair market\nstandard terms.", GREEN),
    ("AI vs AI\nNegotiator", "Two AI agents debate\nyour contract — Advocate\nvs Challenger.", AMBER),
    ("Verified\nFinal Report", "Hallucination guard catches\nerrors before the final\nadvisory is delivered.", ORANGE_LIGHT),
]

for i, (title, desc, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.8) + Inches(row * 2.2)

    # Card
    card = add_rect(s, x, y, Inches(3.7), Inches(1.9), SURFACE_200, border_color=SURFACE_300)
    # Top accent bar
    add_thin_bar(s, x + Inches(0.2), y + Inches(0.18), Inches(0.55), Inches(0.05), accent)
    # Title
    add_text(s, x + Inches(0.2), y + Inches(0.35), Inches(3.3), Inches(0.6),
             title, size=14, color=CREAM, bold=True)
    # Description
    add_text(s, x + Inches(0.2), y + Inches(0.95), Inches(3.3), Inches(0.85),
             desc, size=11, color=CREAM_DIM)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE (DAG Flowchart)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "ARCHITECTURE",
               "Multi-Agent DAG Pipeline — 3 Waves + Verification")

BOX_W = Inches(2.2)
BOX_H = Inches(0.85)

# Wave labels on left
waves = [
    (Inches(0.3), Inches(2.45), "INPUT", CREAM_DIM),
    (Inches(0.3), Inches(3.55), "WAVE 1", ORANGE),
    (Inches(0.3), Inches(4.7), "WAVE 2", ORANGE),
    (Inches(0.3), Inches(5.85), "WAVE 3", ORANGE),
]
for lx, ly, lbl, clr in waves:
    add_text(s, lx, ly, Inches(1), Inches(0.35), lbl, size=10, color=clr, bold=True)
    add_line(s, Inches(1.3), ly + Inches(0.18), Inches(1.8), ly + Inches(0.18), SURFACE_300, 1)

# ── Agent boxes ──
# Upload
ux, uy = Inches(5.2), Inches(2.2)
add_rect(s, ux, uy, BOX_W, BOX_H, SURFACE_200, border_color=SURFACE_300,
         text="Document Ingestion", font_size=12, font_color=CREAM, bold=True)
add_text(s, ux, uy + BOX_H, BOX_W, Inches(0.22),
         "PyMuPDF + OCR + LLM Segmenter", size=8, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Wave 1
w1y = Inches(3.4)
cls_x, sim_x = Inches(3.2), Inches(7.7)

add_rect(s, cls_x, w1y, BOX_W, BOX_H, SURFACE_200, border_color=ORANGE,
         text="Clause Classifier", font_size=12, font_color=CREAM, bold=True)
add_text(s, cls_x, w1y + BOX_H, BOX_W, Inches(0.22),
         "Claude Sonnet 4.5  (blocking)", size=8, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

add_rect(s, sim_x, w1y, BOX_W, BOX_H, SURFACE_200, border_color=SURFACE_300,
         text="Plain Language", font_size=12, font_color=CREAM, bold=True)
add_text(s, sim_x, w1y + BOX_H, BOX_W, Inches(0.22),
         "Claude Sonnet 4.5  (non-blocking)", size=8, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

# "parallel" label between them
add_text(s, Inches(5.6), w1y + Inches(0.25), Inches(1.8), Inches(0.3),
         "parallel", size=9, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Wave 2
w2y = Inches(4.55)
risk_x, bench_x = Inches(3.2), Inches(7.7)

add_rect(s, risk_x, w2y, BOX_W, BOX_H, SURFACE_200, border_color=RED,
         text="Risk Analyzer", font_size=12, font_color=CREAM, bold=True)
add_text(s, risk_x, w2y + BOX_H, BOX_W, Inches(0.22),
         "Claude Sonnet 4.5", size=8, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

add_rect(s, bench_x, w2y, BOX_W, BOX_H, SURFACE_200, border_color=GREEN,
         text="Benchmark", font_size=12, font_color=CREAM, bold=True)
add_text(s, bench_x, w2y + BOX_H, BOX_W, Inches(0.22),
         "Claude Sonnet 4.5", size=8, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

add_text(s, Inches(5.6), w2y + Inches(0.25), Inches(1.8), Inches(0.3),
         "parallel", size=9, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Wave 3
w3y = Inches(5.7)
adv_x, ver_x = Inches(3.2), Inches(7.7)

add_rect(s, adv_x, w3y, BOX_W, BOX_H, SURFACE_200, border_color=AMBER,
         text="Advisory Synthesis", font_size=12, font_color=CREAM, bold=True)
add_text(s, adv_x, w3y + BOX_H, BOX_W, Inches(0.22),
         "Claude Sonnet 4.5", size=8, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

add_rect(s, ver_x, w3y, BOX_W, BOX_H, SURFACE_200, border_color=BLUE_ACCENT,
         text="Hallucination Verifier", font_size=12, font_color=CREAM, bold=True)
add_text(s, ver_x, w3y + BOX_H, BOX_W, Inches(0.22),
         "Claude Sonnet 4.5  (critic)", size=8, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

add_text(s, Inches(5.6), w3y + Inches(0.25), Inches(1.8), Inches(0.3),
         "sequential  \u2192", size=9, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# ── Arrows ──
umid = ux + BOX_W // 2
cmid = cls_x + BOX_W // 2
smid = sim_x + BOX_W // 2

# Upload -> split to W1
add_arrow(s, umid, uy + BOX_H + Inches(0.22), cmid, w1y - Inches(0.05))
add_arrow(s, umid, uy + BOX_H + Inches(0.22), smid, w1y - Inches(0.05))
# W1 -> W2
add_arrow(s, cmid, w1y + BOX_H + Inches(0.22), cmid, w2y - Inches(0.05))
add_arrow(s, smid, w1y + BOX_H + Inches(0.22), smid, w2y - Inches(0.05))
# W2 -> W3
add_arrow(s, cmid, w2y + BOX_H + Inches(0.22), cmid, w3y - Inches(0.05))
add_arrow(s, smid, w2y + BOX_H + Inches(0.22), smid, w3y - Inches(0.05))
# Advisor -> Verifier horizontal
add_arrow(s, adv_x + BOX_W + Inches(0.1), w3y + BOX_H // 2,
          ver_x - Inches(0.1), w3y + BOX_H // 2, BLUE_ACCENT)

# Right side annotations
annotations = [
    ("6 Specialized AI Agents", CREAM),
    ("3 Parallel Execution Waves", CREAM),
    ("Graceful Partial Failure Handling", CREAM),
    ("Cross-Clause Verification", CREAM),
    ("Real-Time SSE Progress Streaming", CREAM),
    ("Multi-Model via OpenRouter", CREAM),
]
add_text(s, Inches(10.5), Inches(2.2), Inches(2.5), Inches(0.3),
         "KEY PROPERTIES", size=10, color=ORANGE, bold=True)
add_line(s, Inches(10.5), Inches(2.5), Inches(12.5), Inches(2.5), ORANGE, 1)
for i, (ann, clr) in enumerate(annotations):
    y = Inches(2.7) + Inches(i * 0.4)
    add_circle(s, Inches(10.5), y + Inches(0.05), Inches(0.15), ORANGE)
    add_text(s, Inches(10.85), y, Inches(2.3), Inches(0.35),
             ann, size=9, color=CREAM_DIM)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "TECH STACK",
               "Built for Speed, Reliability & Multi-Model Intelligence")

tech = [
    ("Backend", "Python 3.11+ / FastAPI\nAsync everywhere\nSSE streaming via Starlette", ORANGE),
    ("AI / LLM", "Claude Sonnet 4.5 (Anthropic)\nGPT-4o-mini / Gemini 2.0 Flash\nRouted via OpenRouter", BLUE_ACCENT),
    ("RAG Pipeline", "sentence-transformers\nChromaDB vector store\nHybrid search + LLM reranking", GREEN),
    ("Document Processing", "PyMuPDF (PDF parsing)\npytesseract (OCR fallback)\nfpdf2 (report generation)", AMBER),
    ("Frontend", "React 19 + TypeScript\nTailwind CSS 4 + Framer Motion\nVite 8 build", ORANGE_LIGHT),
    ("Data & Config", "Pydantic v2 (validation)\nFile-based JSON storage\nSHA-256 deduplication", CREAM_DIM),
]

for i, (cat, items, accent) in enumerate(tech):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.3) + Inches(row * 2.45)

    # Card
    add_rect(s, x, y, Inches(3.7), Inches(2.1), SURFACE_200, border_color=SURFACE_300)
    # Top accent
    add_thin_bar(s, x, y, Inches(3.7), Inches(0.05), accent)
    # Category
    add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(3.2), Inches(0.35),
             cat, size=16, color=accent, bold=True)
    # Line under category
    add_line(s, x + Inches(0.25), y + Inches(0.58), x + Inches(1.8), y + Inches(0.58), SURFACE_300, 1)
    # Items
    add_text(s, x + Inches(0.25), y + Inches(0.7), Inches(3.2), Inches(1.2),
             items, size=12, color=CREAM_DIM)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRODUCT: ONBOARDING & CHAT (image left, text right)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "PRODUCT", "Onboarding & Intelligent Chat")

# Welcome screenshot — left side, full bleed
IMG_W = Inches(7.0)
IMG_H = int(IMG_W * SS_RATIO)
ss8 = os.path.join(SS_DIR, "ss_8.jpg")
add_img(s, ss8, Inches(0.4), Inches(1.7), IMG_W)

# Subtle shadow/border effect behind image
# (placed before image in z-order won't work in pptx easily, so skip)

# Text content — right side
tx = Inches(7.8)
add_text(s, tx, Inches(2.0), Inches(5), Inches(0.4),
         "Personalized Onboarding", size=20, color=CREAM, bold=True)
add_thin_bar(s, tx, Inches(2.5), Inches(2), Inches(0.04), ORANGE)

onboarding_points = [
    ("Location", "Jurisdiction-aware analysis\n(India, USA, UK, Canada, UAE, Singapore)"),
    ("Profession", "Adapts depth of analysis\n(Lawyer, Business Owner, Student...)"),
    ("Use Case", "Prioritizes relevant features\n(Risk, Compliance, Negotiation...)"),
]
for i, (label, desc) in enumerate(onboarding_points):
    y = Inches(2.8) + Inches(i * 1.05)
    add_circle(s, tx, y + Inches(0.02), Inches(0.32), ORANGE, str(i+1), fs=12)
    add_text(s, tx + Inches(0.5), y, Inches(4.2), Inches(0.3),
             label, size=14, color=CREAM, bold=True)
    add_text(s, tx + Inches(0.5), y + Inches(0.32), Inches(4.2), Inches(0.6),
             desc, size=11, color=CREAM_DIM)

# Second screenshot below — chat
ss3 = os.path.join(SS_DIR, "ss_3.jpg")
add_text(s, tx, Inches(5.9), Inches(5), Inches(0.35),
         "Context-aware legal Q&A grounded\nin uploaded documents via RAG.",
         size=11, color=CREAM_DIM)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — PRODUCT: CHAT & LIBRARY (image right, text left)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "PRODUCT", "AI-Powered Legal Chat & Knowledge Base")

# Text on left
tx = Inches(0.8)
add_text(s, tx, Inches(2.0), Inches(4.5), Inches(0.4),
         "Smart Legal Q&A", size=20, color=CREAM, bold=True)
add_thin_bar(s, tx, Inches(2.5), Inches(2), Inches(0.04), ORANGE)

chat_points = [
    "Ask anything about your uploaded document",
    "Answers grounded in actual clauses via RAG",
    "Jurisdiction-aware (Indian legal context)",
    "Conversation memory across questions",
    "Source citations in every response",
]
for i, pt in enumerate(chat_points):
    y = Inches(2.8) + Inches(i * 0.48)
    add_circle(s, tx, y + Inches(0.05), Inches(0.18), ORANGE)
    add_text(s, tx + Inches(0.35), y, Inches(4.2), Inches(0.4),
             pt, size=12, color=CREAM_DIM)

add_text(s, tx, Inches(5.3), Inches(4.5), Inches(0.4),
         "Legal Knowledge Library", size=18, color=CREAM, bold=True)
add_thin_bar(s, tx, Inches(5.75), Inches(2), Inches(0.04), GREEN)
add_text(s, tx, Inches(5.95), Inches(4.5), Inches(0.8),
         "Built-in Indian legal references — Contract Act 1872,\n"
         "Rent Control Acts, Shops & Establishments Act, and more.",
         size=11, color=CREAM_DIM)

# Chat screenshot — right
ss1 = os.path.join(SS_DIR, "ss_3.jpg")
IMG_W2 = Inches(6.8)
add_img(s, ss1, Inches(5.8), Inches(1.7), IMG_W2)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRODUCT: RISK SCORE & PIPELINE & NEGOTIATOR
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "PRODUCT", "Analysis Tools — Risk, Pipeline & Negotiation")

# Risk Score — large, left side
ss11 = os.path.join(SS_DIR, "ss_11.jpg")
IMG_W3 = Inches(6.5)
add_img(s, ss11, Inches(0.3), Inches(1.7), IMG_W3)
add_text(s, Inches(0.3), Inches(1.7) + int(IMG_W3 * SS_RATIO) + Inches(0.08),
         IMG_W3, Inches(0.3),
         "Risk Score Dashboard — Clause-Level Breakdown with Verdict",
         size=10, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Right side — stacked: pipeline on top, negotiator below
right_x = Inches(7.1)
IMG_W4 = Inches(5.8)
IMG_H4 = int(IMG_W4 * SS_RATIO)

# Pipeline
ss4 = os.path.join(SS_DIR, "ss_4.jpg")
add_img(s, ss4, right_x, Inches(1.7), IMG_W4)
add_text(s, right_x, Inches(1.7) + IMG_H4 + Inches(0.05),
         IMG_W4, Inches(0.25),
         "Real-Time Agent Pipeline Progress (SSE)",
         size=9, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Negotiator
ss5 = os.path.join(SS_DIR, "ss_5.jpg")
neg_y = Inches(1.7) + IMG_H4 + Inches(0.35)
add_img(s, ss5, right_x, neg_y, IMG_W4)
add_text(s, right_x, neg_y + IMG_H4 + Inches(0.05),
         IMG_W4, Inches(0.25),
         "AI vs AI Negotiator — Advocate vs Challenger Debate",
         size=9, color=CREAM_DIM, align=PP_ALIGN.CENTER)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — KEY DIFFERENTIATORS
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
section_header(s, "WHAT SETS US APART",
               "Beyond Summarization — A True Legal Action Agent")

diffs = [
    ("Multi-Model Intelligence",
     "Each agent uses the optimal AI model for its task — Claude for deep analysis, "
     "GPT-4o-mini for chat, Gemini for creative negotiation. Not one model fits all.",
     ORANGE),
    ("Hallucination Guard",
     "Every clause reference is cross-checked against source. The Verifier catches "
     "factual errors, cross-clause conflicts, and fabricated citations before the final report.",
     RED),
    ("Jurisdiction-Aware",
     "Built-in Indian legal context — Rent Control Acts, Section 27, Shops & Establishments. "
     "Not generic AI summarization, but legally grounded analysis.",
     GREEN),
    ("Real-Time Transparency",
     "Users see exactly which agent is running, on which model, in real-time via SSE streaming. "
     "Full pipeline visibility — no black box.",
     BLUE_ACCENT),
    ("Actionable Decisions",
     "Delivers a verdict (Sign / Negotiate / Walk Away), negotiation playbook, and prioritized "
     "action items — not just a summary.",
     AMBER),
]

for i, (title, desc, accent) in enumerate(diffs):
    y = Inches(2.2) + Inches(i * 1.02)

    # Left accent + number
    add_circle(s, Inches(0.8), y + Inches(0.05), Inches(0.32), accent, str(i+1), fs=13)

    # Title
    add_text(s, Inches(1.35), y, Inches(3), Inches(0.35),
             title, size=15, color=CREAM, bold=True)
    # Separator
    add_text(s, Inches(4.2), y + Inches(0.05), Inches(0.3), Inches(0.3),
             "|", size=14, color=SURFACE_300)
    # Description
    add_text(s, Inches(4.5), y + Inches(0.02), Inches(8), Inches(0.5),
             desc, size=11, color=CREAM_DIM)

    # Subtle line under each
    if i < len(diffs) - 1:
        add_line(s, Inches(0.8), y + Inches(0.85), Inches(12.5), y + Inches(0.85),
                 SURFACE_300, 1)

slide_footer(s)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

# Top accent
add_thin_bar(s, Inches(0), Inches(0), W, Inches(0.05), ORANGE)

# Logo
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.92), Inches(1.5), Inches(1.5))

# Thank you
add_text(s, Inches(2), Inches(3.1), Inches(9.333), Inches(0.9),
         "Thank You", size=52, color=CREAM, bold=True, align=PP_ALIGN.CENTER)

# Quote from problem statement
add_text(s, Inches(2.5), Inches(4.2), Inches(8.333), Inches(0.8),
         '"Making legal understanding accessible,\nsafe, and actionable for everyone."',
         size=18, color=ORANGE_SOFT, align=PP_ALIGN.CENTER)

# Divider
add_line(s, Inches(5.5), Inches(5.3), Inches(7.833), Inches(5.3), SURFACE_300, 1)

# Subtitle
add_text(s, Inches(2.5), Inches(5.5), Inches(8.333), Inches(0.5),
         "Suits AI  —  Multi-Agent Legal Document Intelligence",
         size=13, color=CREAM_DIM, align=PP_ALIGN.CENTER)

# Bottom bar
add_thin_bar(s, Inches(0), Inches(7.1), W, Inches(0.4), SURFACE_200)
add_text(s, Inches(0.8), Inches(7.13), Inches(5), Inches(0.35),
         "RNSIT Agentic AI Hackathon  |  Problem Statement 3",
         size=9, color=CREAM_DIM)
add_text(s, Inches(8), Inches(7.13), Inches(4.5), Inches(0.35),
         "AI Legal Document Action Agent", size=9, color=ORANGE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
